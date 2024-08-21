import urllib.request
import urllib.error
import json
from docx import Document
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv
from messages import prompt_prefix1, prompt_prefix2

load_dotenv()

def makeApiCall(apiKey, content):
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        'x-api-key': apiKey,
        'anthropic-version': '2023-06-01',
        'content-type': 'application/json'
    }
    data = json.dumps({
        "model": "claude-3-sonnet-20240229",
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": content}]
    }).encode('utf-8')

    print(f"Attempting to access: {url}")
    print(f"With headers: {headers}")

    req = urllib.request.Request(url, data=data, headers=headers, method='POST')

    try:
        with urllib.request.urlopen(req) as response:
            output = response.read()
            return json.loads(output)['content'][0]['text']
    except urllib.error.HTTPError as e:
        print(f"HTTP error: {e.code} - {e.reason}")
        print(f"Response body: {e.read().decode()}")
    except urllib.error.URLError as e:
        print(f"URL error: {e}")
    except Exception as e:
        print(f"General error: {e}")
    return None

def extractTopicLines(output):
    extracted_dict = []
    lines = output.split("\n")
    for i, line in enumerate(lines):
        if line.startswith("**") and line.endswith("**"):
            key_topic = line.strip("**")
            starting_lines = "\n".join(lines[i+1:i+3])
            extracted_dict.append({"key_topic": key_topic, "starting_lines": starting_lines})
    return extracted_dict

def saveContentToFile(filePath, content):
    with open(filePath, 'w', encoding='utf-8') as file:
        file.write(content + "\n")

def appendContentToFile(filePath, keyTopic, content):
    with open(filePath, 'a', encoding='utf-8') as file:
        file.write(f"**{keyTopic}**\n{content}\n")

def read_document(file_path):
    _, file_extension = os.path.splitext(file_path)
    if file_extension.lower() == '.docx':
        return read_word_document(file_path)
    else:
        return read_text_document(file_path)

def read_word_document(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    saveContentToFile(wordFilePath, '\n'.join(full_text))
    return '\n'.join(full_text)

def read_text_document(file_path):
    with open(file_path, 'r', encoding="utf-8") as file:
        return file.read()

def parse_input(filename):
    with open(filename, 'r', encoding='utf-8') as file:
        content = file.read()
    
    sections = re.split(r'\*\*(.*?)\*\*', content)[1:]
    parsed_data = []
    i = 0
    while i < len(sections):
        topic = sections[i].strip()
        points = [p.strip() for p in sections[i+1].split('*') if p.strip()]
        parsed_data.append((topic, points))
        i += 2
    return parsed_data

def add_slide(prs, layout, title, points):
    slide = prs.slides.add_slide(layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    if points and slide.shapes.placeholders[1].has_text_frame:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = ""
        
        for point in points:
            # print("POINT HERE")
            # print(point)
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    
    return slide

def update_presentation(parsed_data, pptx_path):
    prs = Presentation()  # Start a new presentation instead of loading an existing one
    
    title_layout = None
    title_content_layout = None
    for layout in prs.slide_layouts:
        if layout.name == 'Title Slide':
            title_layout = layout
        elif layout.name == 'Title and Content':
            title_content_layout = layout
    
    if not title_layout:
        title_layout = prs.slide_layouts[0]
    if not title_content_layout:
        title_content_layout = prs.slide_layouts[1]
    # print("#######################")
    # print(parsed_data)
    for topic, points in parsed_data:
        if not points:
            add_slide(prs, title_layout, topic, [])
        else:
            # for point in points:
            #     print(point)
            add_slide(prs, title_content_layout, topic, points)
                # points = points[7:]
    
    prs.save(pptx_path)  # Save the new presentation

def main(apiKey, articlePath, prompt1OutputPath, powerPointStructurePath, pptxPath):
    # Create empty files at the start of the session
    open(prompt1OutputPath, 'w').close()
    open(powerPointStructurePath, 'w').close()
    
    # read the original doc
    document = read_document(articlePath)
    prompt1_content = prompt_prefix1 + document
    output1 = makeApiCall(apiKey, prompt1_content)
    promptTemplate2 = prompt_prefix2

    if output1:
        saveContentToFile(prompt1OutputPath, output1)
        extracted_dict = extractTopicLines(output1)
    
        for i, entry in enumerate(extracted_dict):
            keyTopic = entry["key_topic"]
            startingLines = entry["starting_lines"]
            nextTopicStartingLines = extracted_dict[i + 1]["starting_lines"] if i + 1 < len(extracted_dict) else ""

            prompt2_content = promptTemplate2.replace("{{keytopic}}", keyTopic).replace("{{startingLines}}", startingLines).replace("{{nextTopicStartingLines}}", nextTopicStartingLines).replace("<paste full document here>", document)
            output2 = makeApiCall(apiKey, prompt2_content)

            if output2:
                appendContentToFile(powerPointStructurePath, keyTopic, output2)

    try:
        parsed_data = parse_input(powerPointStructurePath)
        # print(parsed_data)
        if not parsed_data:
            print("Error: No topics found in the input file.")
            return
        
        update_presentation(parsed_data, pptxPath)
        print(f"Slides added to the PowerPoint presentation successfully: {pptxPath}")
    except FileNotFoundError:
        print("Error: input.txt file or specified PowerPoint file not found.")
    except Exception as e:
        print(f"An error occurred: {str(e)}")

if __name__ == "__main__":
    apiKey = os.getenv('ANTHROPIC_API_KEY') 
    articlePath = "doc.txt"
    prompt1OutputPath = "prompt1Output.txt"
    powerPointStructurePath = "powerPointStructure.txt"
    wordFilePath = "wordFile.txt"
    pptxPath = "test.pptx"

    main(apiKey, articlePath, prompt1OutputPath, powerPointStructurePath, pptxPath)
